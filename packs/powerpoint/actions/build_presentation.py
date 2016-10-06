import os
from pptx import Presentation

from st2actions.runners.pythonrunner import Action
from st2common.exceptions.action import InvalidActionParameterException

__all__ = [
    'BuildPresentationAction',
]


class BuildPresentationAction(Action):
    def __init__(self, config):
        super(BuildPresentationAction, self).__init__(config)
        self.template = self.config['template']

    def run(self, path, slides):
        if self.template is not None and self.template != '':
            if os.path.exists(self.template):
                prs = Presentation(self.template)
        else:
            prs = Presentation()

        if not isinstance(slides, list):
            raise InvalidActionParameterException("slides must be a list")

        for slide in slides:
            slide = prs.slides.add_slide(slide.get('layout', 1))
            shapes = slide.shapes

            title_shape = shapes.title
            body_shape = shapes.placeholders[1]

            title_shape.text = slide.get('title', '')

            tf = body_shape.text_frame
            tf.text = slide.get('text', '')

        result = prs.save(path)
        return result
